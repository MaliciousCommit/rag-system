# ingestion/loaders.py

"""
Document Loaders — Extract raw text from any file type

SUPPORTED FORMATS:
    PDF  → pypdf (text-based) or Google Doc AI (scanned/OCR)
    DOCX → python-docx
    PPTX → python-pptx
    HTML → BeautifulSoup4

WHY SEPARATE LOADERS PER TYPE?
    Each format stores content differently:
    - PDFs have pages, fonts, layout metadata
    - DOCX has paragraphs, headings, tables
    - PPTX has slides, shapes, speaker notes
    - HTML has tags, scripts, navigation to strip out

    One generic loader would do a poor job on all of them.
    Specialized loaders extract clean, structured text.

OUTPUT CONTRACT:
    Every loader returns List[Document] where each Document has:
    - page_content : str   → the actual text
    - metadata     : dict  → source, page/slide number, file type etc.

    Consistent output means the Chunker doesn't care which
    loader produced the document — it works the same way for all.
"""

import os
import logging
from pathlib import Path
from typing import List
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


# ─── Document Model ───────────────────────────────────────────────────────────


@dataclass
class Document:
    """
    Represents a single unit of extracted content.

    WHY a custom Document class instead of LangChain's?
        More control over metadata fields.
        Easier to serialize to JSON for GCS storage.
        LangChain Document is used downstream in the chunker.
    """

    page_content: str  # The actual text content
    metadata: dict = field(default_factory=dict)
    # metadata keys we always populate:
    #   source    : original file path or URL
    #   file_type : "pdf" | "docx" | "pptx" | "html"
    #   page      : page/slide number (0-indexed)
    #   total_pages: total pages in document


# ─── PDF Loader ───────────────────────────────────────────────────────────────


def load_pdf(file_path: str) -> List[Document]:
    """
    Extract text from a PDF file page by page.

    WHY page by page (not whole document)?
        Preserves page metadata — useful for citations.
        "This answer comes from page 5 of policy.pdf" is more
        useful than "this comes from policy.pdf".
        Also prevents huge single chunks from very long PDFs.

    LIMITATION:
        pypdf only works on text-based PDFs.
        Scanned PDFs (image-based) need OCR → use load_pdf_ocr() instead.
        We auto-detect which to use in the main pipeline.
    """
    try:
        import pypdf

        documents = []
        path = Path(file_path)

        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            total_pages = len(reader.pages)

            logger.info(f"[PDF Loader] Reading {path.name} ({total_pages} pages)")

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()

                # Skip empty pages — common in PDFs with full-page images
                if not text or not text.strip():
                    logger.debug(f"[PDF Loader] Skipping empty page {page_num + 1}")
                    continue

                documents.append(
                    Document(
                        page_content=text.strip(),
                        metadata={
                            "source": str(path.name),
                            "file_path": str(file_path),
                            "file_type": "pdf",
                            "page": page_num + 1,  # 1-indexed for humans
                            "total_pages": total_pages,
                        },
                    )
                )

        logger.info(f"[PDF Loader] Extracted {len(documents)} pages from {path.name}")
        return documents

    except Exception as e:
        logger.error(f"[PDF Loader] Failed on {file_path}: {e}")
        return []


# ─── DOCX Loader ──────────────────────────────────────────────────────────────


def load_docx(file_path: str) -> List[Document]:
    """
    Extract text from a Word document (.docx).

    STRATEGY:
        Group paragraphs into logical sections using headings.
        A heading + its following paragraphs = one Document.
        This preserves semantic structure better than fixed splits.

    WHY group by heading?
        "Section 3: Refund Policy" + its paragraphs is one logical unit.
        If we split mid-section, the retriever might fetch half an answer.
        Grouping by heading gives the LLM complete context per chunk.
    """
    try:
        from docx import Document as DocxDocument

        path = Path(file_path)
        doc = DocxDocument(file_path)

        logger.info(f"[DOCX Loader] Reading {path.name}")

        documents = []
        current_section = []
        current_heading = "Introduction"
        section_num = 0

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # Detect headings by style name
            style_name = (para.style.name or "") if para.style else ""
            is_heading = style_name.startswith("Heading")

            if is_heading and current_section:
                # Save the completed section as a Document
                documents.append(
                    Document(
                        page_content="\n".join(current_section),
                        metadata={
                            "source": str(path.name),
                            "file_path": str(file_path),
                            "file_type": "docx",
                            "section": current_heading,
                            "page": section_num,
                            "total_pages": -1,  # DOCX has no page concept
                        },
                    )
                )
                current_section = []
                section_num += 1
                current_heading = text
            else:
                current_section.append(text)

        # Don't forget the last section
        if current_section:
            documents.append(
                Document(
                    page_content="\n".join(current_section),
                    metadata={
                        "source": str(path.name),
                        "file_path": str(file_path),
                        "file_type": "docx",
                        "section": current_heading,
                        "page": section_num,
                        "total_pages": section_num,
                    },
                )
            )

        logger.info(
            f"[DOCX Loader] Extracted {len(documents)} sections from {path.name}"
        )
        return documents

    except Exception as e:
        logger.error(f"[DOCX Loader] Failed on {file_path}: {e}")
        return []


# ─── PPTX Loader ──────────────────────────────────────────────────────────────


def load_pptx(file_path: str) -> List[Document]:
    """
    Extract text from a PowerPoint presentation (.pptx).
    One Document per slide — includes title, body, speaker notes.
    """
    try:
        from pptx import Presentation

        path = Path(file_path)
        prs = Presentation(file_path)
        total_slides = len(prs.slides)

        logger.info(f"[PPTX Loader] Reading {path.name} ({total_slides} slides)")

        documents = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            title = ""

            for shape in slide.shapes:
                # Skip shapes with no text
                if not shape.has_text_frame:
                    continue

                # ✅ FIX: getattr safely accesses text_frame
                # BaseShape doesn't declare text_frame but subclasses do.
                # getattr bypasses Pylance's type check while staying safe at runtime.
                text_frame = getattr(shape, "text_frame", None)
                if text_frame is None:
                    continue

                for para in text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue

                    # Detect title shape by name
                    if "title" in shape.name.lower():
                        title = text
                    else:
                        slide_texts.append(text)

            # Add speaker notes if they exist
            if slide.has_notes_slide:
                notes_frame = getattr(slide.notes_slide, "notes_text_frame", None)
                if notes_frame:
                    notes = notes_frame.text.strip()
                    if notes:
                        slide_texts.append(f"[Speaker Notes]: {notes}")

            # Only add slide if it has content
            full_text = f"{title}\n{chr(10).join(slide_texts)}".strip()
            if full_text:
                documents.append(
                    Document(
                        page_content=full_text,
                        metadata={
                            "source": str(path.name),
                            "file_path": str(file_path),
                            "file_type": "pptx",
                            "slide_title": title,
                            "page": slide_num,
                            "total_pages": total_slides,
                        },
                    )
                )

        logger.info(f"[PPTX Loader] Extracted {len(documents)} slides from {path.name}")
        return documents

    except Exception as e:
        logger.error(f"[PPTX Loader] Failed on {file_path}: {e}")
        return []


# ─── HTML Loader ──────────────────────────────────────────────────────────────


def load_html(file_path: str) -> List[Document]:
    """
    Extract clean text from an HTML file.

    STRATEGY:
        Use BeautifulSoup to strip all tags.
        Remove scripts, styles, navigation — keep main content only.
        Group by <h1>/<h2> headings into sections.

    WHY BeautifulSoup over a simple regex strip?
        HTML is nested and irregular. Regex breaks on edge cases.
        BS4 parses the DOM properly and handles malformed HTML.
    """
    try:
        from bs4 import BeautifulSoup

        path = Path(file_path)

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()

        logger.info(f"[HTML Loader] Reading {path.name}")

        soup = BeautifulSoup(html_content, "lxml")

        # Remove noise elements — these never contain useful RAG content
        for tag in soup(
            ["script", "style", "nav", "footer", "header", "aside", "form", "iframe"]
        ):
            tag.decompose()

        # Extract sections grouped by headings
        documents = []
        current_section = []
        current_heading = "Main Content"
        section_num = 0

        for element in soup.find_all(["h1", "h2", "h3", "p", "li", "td"]):
            text = element.get_text(separator=" ", strip=True)
            if not text:
                continue

            if element.name in ["h1", "h2", "h3"]:
                if current_section:
                    documents.append(
                        Document(
                            page_content="\n".join(current_section),
                            metadata={
                                "source": str(path.name),
                                "file_path": str(file_path),
                                "file_type": "html",
                                "section": current_heading,
                                "page": section_num,
                            },
                        )
                    )
                    current_section = []
                    section_num += 1
                current_heading = text
            else:
                current_section.append(text)

        # Last section
        if current_section:
            documents.append(
                Document(
                    page_content="\n".join(current_section),
                    metadata={
                        "source": str(path.name),
                        "file_path": str(file_path),
                        "file_type": "html",
                        "section": current_heading,
                        "page": section_num,
                    },
                )
            )

        logger.info(
            f"[HTML Loader] Extracted {len(documents)} sections from {path.name}"
        )
        return documents

    except Exception as e:
        logger.error(f"[HTML Loader] Failed on {file_path}: {e}")
        return []


# ─── Router — Auto-detect file type ───────────────────────────────────────────


def load_document(file_path: str) -> List[Document]:
    """
    Auto-detect file type and route to the correct loader.

    This is the ONLY function the pipeline calls externally.
    It hides the complexity of multiple loaders behind one interface.

    Args:
        file_path: Path to any supported document

    Returns:
        List[Document] regardless of file type
    """
    ext = Path(file_path).suffix.lower()

    loader_map = {
        ".pdf": load_pdf,
        ".docx": load_docx,
        ".doc": load_docx,
        ".pptx": load_pptx,
        ".ppt": load_pptx,
        ".html": load_html,
        ".htm": load_html,
    }

    loader_fn = loader_map.get(ext)

    if loader_fn is None:
        logger.warning(f"[Loader] Unsupported file type: {ext} — skipping {file_path}")
        return []

    logger.info(f"[Loader] Routing {Path(file_path).name} → {loader_fn.__name__}")
    return loader_fn(file_path)
